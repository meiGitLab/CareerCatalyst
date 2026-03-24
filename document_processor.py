# document_processor.py

import os
from typing import List
import pandas as pd
from langchain.docstore.document import Document as LangchainDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Alternative PDF processing without unstructured
import PyPDF2
import pdfplumber

class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
    
    def process_pdf(self, file_path: str) -> List[LangchainDocument]:
        """Process PDF files using PyPDF2 and pdfplumber for text and tables"""
        documents = []
        
        try:
            # Method 1: Extract text with PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        documents.append(LangchainDocument(
                            page_content=text,
                            metadata={
                                "source": file_path, 
                                "type": "pdf", 
                                "filename": os.path.basename(file_path),
                                "page": page_num + 1
                            }
                        ))
            
            # Method 2: Extract tables with pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        tables = page.extract_tables()
                        for table_num, table in enumerate(tables):
                            if table:
                                table_text = "\n".join(["\t".join(map(str, row)) for row in table])
                                documents.append(LangchainDocument(
                                    page_content=f"Table {table_num + 1}:\n{table_text}",
                                    metadata={
                                        "source": file_path, 
                                        "type": "pdf_table", 
                                        "filename": os.path.basename(file_path),
                                        "page": page_num + 1,
                                        "table": table_num + 1
                                    }
                                ))
            except Exception as e:
                print(f"PDF table extraction error {file_path}: {e}")
                
        except Exception as e:
            print(f"Error processing PDF file {file_path}: {e}")
        
        return documents
    
    def process_docx(self, file_path: str) -> List[LangchainDocument]:
        """Process DOCX files"""
        try:
            import docx
            doc = docx.Document(file_path)
            content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
            
            documents = [LangchainDocument(
                page_content=content,
                metadata={"source": file_path, "type": "docx", "filename": os.path.basename(file_path)}
            )]
            
            return self.text_splitter.split_documents(documents)
        except Exception as e:
            print(f"Error processing DOCX file {file_path}: {e}")
            return []
    
    def process_excel(self, file_path: str) -> List[LangchainDocument]:
        """Process Excel files, extract all worksheet data"""
        documents = []
        
        try:
            # Read all worksheets
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert DataFrame to readable text format
                sheet_content = f"Worksheet: {sheet_name}\n\n"
                sheet_content += df.to_string(index=False)
                
                documents.append(LangchainDocument(
                    page_content=sheet_content,
                    metadata={"source": file_path, "type": "excel", "sheet": sheet_name, "filename": os.path.basename(file_path)}
                ))
        except Exception as e:
            print(f"Error processing Excel file {file_path}: {e}")
        
        return documents
    
    def process_ppt(self, file_path: str) -> List[LangchainDocument]:
        """Process PPT files using python-pptx"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                
                if slide_content:
                    content.append(f"Slide {slide_num + 1}:\n" + "\n".join(slide_content))
            
            full_content = "\n\n".join(content)
            
            documents = [LangchainDocument(
                page_content=full_content,
                metadata={"source": file_path, "type": "ppt", "filename": os.path.basename(file_path)}
            )]
            
            return self.text_splitter.split_documents(documents)
        except Exception as e:
            print(f"Error processing PPT file {file_path}: {e}")
            return []
    
    def process_text_file(self, file_path: str) -> List[LangchainDocument]:
        """Process text files (TXT, MD, etc.)"""
        try:
            # Detect file encoding
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            content = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                print(f"Cannot decode file: {file_path}")
                return []
            
            documents = [LangchainDocument(
                page_content=content,
                metadata={"source": file_path, "type": "text", "filename": os.path.basename(file_path)}
            )]
            
            return self.text_splitter.split_documents(documents)
        except Exception as e:
            print(f"Error processing text file {file_path}: {e}")
            return []
    
    def process_csv(self, file_path: str) -> List[LangchainDocument]:
        """Process CSV files"""
        try:
            df = pd.read_csv(file_path)
            content = f"CSV File: {os.path.basename(file_path)}\n\n"
            content += df.to_string(index=False)
            
            documents = [LangchainDocument(
                page_content=content,
                metadata={"source": file_path, "type": "csv", "filename": os.path.basename(file_path)}
            )]
            
            return documents
        except Exception as e:
            print(f"Error processing CSV file {file_path}: {e}")
            return []
    
    def process_directory(self, data_directory: str) -> List[LangchainDocument]:
        """Process all supported files in the directory"""
        all_documents = []
        
        file_processors = {
            '.pdf': self.process_pdf,
            '.docx': self.process_docx,
            '.xlsx': self.process_excel,
            '.xls': self.process_excel,
            '.csv': self.process_csv,
            '.ppt': self.process_ppt,
            '.pptx': self.process_ppt,
            '.txt': self.process_text_file,
            '.md': self.process_text_file,
            '.html': self.process_text_file,
            '.htm': self.process_text_file
        }
        
        if not os.path.exists(data_directory):
            print(f"Data directory does not exist: {data_directory}")
            return all_documents
        
        processed_count = 0
        for root, _, files in os.walk(data_directory):
            for file in files:
                file_path = os.path.join(root, file)
                file_ext = os.path.splitext(file_path)[1].lower()
                
                if file_ext in file_processors:
                    try:
                        print(f"Processing: {file_path}")
                        documents = file_processors[file_ext](file_path)
                        if documents:
                            all_documents.extend(documents)
                            processed_count += 1
                            print(f"  Successfully processed {file}, generated {len(documents)} document chunks")
                        else:
                            print(f"  Processing {file} generated no content")
                    except Exception as e:
                        print(f"Error processing file {file}: {e}")
        
        print(f"\nProcessing completed! Processed {processed_count} files, generated {len(all_documents)} document chunks")
        return all_documents

# Global instance
document_processor = DocumentProcessor()